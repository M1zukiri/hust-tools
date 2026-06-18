"""
PPTX/PPT 转 Markdown 转换器
使用 python-pptx 库
"""

from pathlib import Path
from typing import Optional, Dict, Any, List
import re

from converter_base import BaseConverter, ConversionResult, ConversionStatus


class PptxConverter(BaseConverter):
    """PowerPoint演示文稿转Markdown转换器"""
    
    supported_extensions = ['pptx', 'ppt']
    supported_mime_types = [
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'application/vnd.ms-powerpoint'
    ]
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        self.include_slide_numbers = config.get('include_slide_numbers', True) if config else True
        self.include_notes = config.get('include_notes', False) if config else False
        self.slide_separator = config.get('slide_separator', '---') if config else '---'
    
    def convert(self, file_path: Path, **options) -> ConversionResult:
        """
        将PPTX文件转换为Markdown
        
        Args:
            file_path: PPTX文件路径
            **options: 额外选项
                - include_slide_numbers: 是否包含幻灯片编号
                - include_notes: 是否包含演讲者备注
                - slide_separator: 幻灯片分隔符
                
        Returns:
            ConversionResult: 转换结果
        """
        if not self.validate_file(file_path):
            return self._create_error_result(file_path, "文件验证失败", "pptx")
        
        try:
            from pptx import Presentation
        except ImportError:
            return self._create_error_result(
                file_path,
                "缺少依赖: python-pptx。请运行: pip install python-pptx",
                "pptx"
            )
        
        try:
            prs = Presentation(str(file_path))
            md_content = self._presentation_to_markdown(prs, file_path, **options)
            
            metadata = self.extract_metadata(file_path)
            metadata['slide_count'] = len(prs.slides)
            
            return ConversionResult(
                content=md_content,
                status=ConversionStatus.SUCCESS,
                source_format='pptx',
                metadata=metadata,
                errors=[]
            )
            
        except Exception as e:
            self.logger.error(f"转换失败: {e}")
            return self._create_error_result(file_path, str(e), "pptx")
    
    def _presentation_to_markdown(self, prs, file_path: Path, **options) -> str:
        """
        将Presentation对象转换为Markdown
        
        Args:
            prs: python-pptx Presentation对象
            file_path: 文件路径
            **options: 额外选项
            
        Returns:
            str: Markdown内容
        """
        include_slide_numbers = options.get('include_slide_numbers', self.include_slide_numbers)
        include_notes = options.get('include_notes', self.include_notes)
        slide_separator = options.get('slide_separator', self.slide_separator)
        
        md_parts = []
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_md = self._slide_to_markdown(slide, idx, include_slide_numbers)
            if slide_md.strip():
                md_parts.append(slide_md)
                
                # 添加演讲者备注
                if include_notes and slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        md_parts.append(f"\n**备注:** {notes}\n")
                
                # 添加分隔符
                if idx < len(prs.slides):
                    md_parts.append(f"\n{slide_separator}\n")
        
        return '\n'.join(md_parts)
    
    def _slide_to_markdown(self, slide, slide_number: int, include_number: bool) -> str:
        """
        将单个幻灯片转换为Markdown
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            include_number: 是否包含编号
            
        Returns:
            str: Markdown文本
        """
        md_parts = []
        
        # 幻灯片标题
        if include_number:
            md_parts.append(f"## 幻灯片 {slide_number}\n")
        
        # 提取文本内容
        texts = self._extract_slide_texts(slide)
        
        # 处理标题
        if texts['title']:
            md_parts.append(f"### {texts['title']}\n")
        
        # 处理副标题
        if texts['subtitle']:
            md_parts.append(f"*{texts['subtitle']}*\n")
        
        # 处理正文内容
        for content in texts['contents']:
            if content.strip():
                # 检测列表项
                if content.strip().startswith(('•', '○', '▪', '-', '·')):
                    md_parts.append(f"- {content.strip()[1:].strip()}")
                elif re.match(r'^\d+[.、]', content.strip()):
                    md_parts.append(content.strip())
                else:
                    md_parts.append(content.strip())
        
        return '\n\n'.join(md_parts)
    
    def _extract_slide_texts(self, slide) -> Dict[str, Any]:
        """
        提取幻灯片中的各类文本
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            Dict: 分类的文本内容
        """
        result = {
            'title': '',
            'subtitle': '',
            'contents': []
        }
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text_frame.text.strip()
            if not text:
                continue
            
            # 根据形状类型和位置判断内容类型
            shape_type = str(shape.shape_type)
            
            # 标题通常位于顶部且字体较大
            if 'TITLE' in shape_type or (hasattr(shape, 'top') and shape.top < 1000000):
                if not result['title']:
                    result['title'] = text
                elif not result['subtitle']:
                    result['subtitle'] = text
                else:
                    result['contents'].append(text)
            else:
                # 处理文本框中的段落
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        result['contents'].append(para_text)
        
        return result
    
    def _extract_tables_from_slide(self, slide) -> List[str]:
        """
        从幻灯片中提取表格
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            List[str]: Markdown表格列表
        """
        tables = []
        
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                md_table = self._table_to_markdown(table)
                if md_table:
                    tables.append(md_table)
        
        return tables
    
    def _table_to_markdown(self, table) -> str:
        """
        将表格转换为Markdown
        
        Args:
            table: 表格对象
            
        Returns:
            str: Markdown表格
        """
        if not table.rows:
            return ""
        
        md_lines = []
        rows = list(table.rows)
        
        # 表头
        header_cells = [cell.text.strip() for cell in rows[0].cells]
        md_lines.append('| ' + ' | '.join(header_cells) + ' |')
        md_lines.append('|' + '|'.join(['---' for _ in header_cells]) + '|')
        
        # 数据行
        for row in rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            md_lines.append('| ' + ' | '.join(cells) + ' |')
        
        return '\n'.join(md_lines)
